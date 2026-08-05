import io
import logging
import os
import zipfile

import PyPDF2
import docx
import pandas as pd
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_tesseract_cmd = os.getenv("TESSERACT_CMD", "tesseract")
pytesseract.pytesseract.tesseract_cmd = _tesseract_cmd

logger = logging.getLogger(__name__)


def _clean_ocr_text(text):
    if not text:
        return ""

    cleaned = " ".join(text.split()).strip()
    if len(cleaned) < 20:
        return ""

    alnum_ratio = sum(c.isalnum() for c in cleaned) / max(1, len(cleaned))
    if alnum_ratio < 0.4:
        return ""

    return cleaned


def _ocr_image(image, source=None):
    try:
        raw_text = pytesseract.image_to_string(image)
        return _clean_ocr_text(raw_text)
    except Exception:
        logger.exception("OCR failed for %s", source or "image")
        return ""


def extract_text(file_path):
    if not file_path or not os.path.exists(file_path):
        raise ValueError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        if ext == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                pages = None

                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text += page_text.strip() + "\n"
                        continue

                    try:
                        if pages is None:
                            pages = convert_from_path(file_path)

                        image = pages[i]
                        ocr_text = _ocr_image(image, source=f"{file_path} page {i + 1}")
                        if ocr_text:
                            text += ocr_text + "\n"
                    except Exception:
                        logger.exception("PDF OCR failed for %s page %s", file_path, i + 1)
                        continue

        elif ext == ".docx":
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

            for rel in doc.part.rels.values():
                target_ref = getattr(rel, "target_ref", "")
                if not target_ref or "image" not in target_ref.lower():
                    continue

                image_part = getattr(rel, "target_part", None)
                if not image_part:
                    continue

                try:
                    image = Image.open(io.BytesIO(image_part.blob))
                    ocr_text = _ocr_image(image, source=f"{file_path} docx image")
                    if ocr_text:
                        text += "\n" + ocr_text + "\n"
                except Exception:
                    logger.exception("DOCX image OCR failed for %s", file_path)
                    continue

        elif ext == ".pptx":
            try:
                prs = Presentation(file_path)
            except (zipfile.BadZipFile, IOError, ValueError) as exc:
                logger.exception("Invalid PPTX file %s", file_path)
                raise ValueError("Uploaded PPTX appears to be corrupted or unreadable.") from exc

            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text += cell.text + "\n"
                    elif hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = Image.open(io.BytesIO(shape.image.blob))
                            ocr_text = _ocr_image(image, source=f"{file_path} pptx image")
                            if ocr_text:
                                text += ocr_text + "\n"
                        except Exception:
                            logger.exception("PPTX image OCR failed for %s", file_path)
                            continue

        elif ext == ".xlsx":
            try:
                df = pd.read_excel(file_path, nrows=1000, dtype=str)
                text = df.to_string(index=False)
            except Exception:
                logger.exception("Failed to read XLSX file %s", file_path)
                raise ValueError("Uploaded XLSX file could not be read.")

        elif ext == ".csv":
            try:
                chunks = pd.read_csv(file_path, dtype=str, low_memory=False, chunksize=1000)
                first_chunks = [chunk for _, chunk in zip(range(3), chunks)]
                if first_chunks:
                    df = pd.concat(first_chunks, ignore_index=True)
                    text = df.to_string(index=False)
                else:
                    text = ""
            except Exception:
                logger.exception("Failed to read CSV file %s", file_path)
                raise ValueError("Uploaded CSV file could not be read.")

        elif ext == ".txt":
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except Exception:
                logger.exception("Failed to read TXT file %s", file_path)
                raise ValueError("Uploaded TXT file could not be read.")

        else:
            raise ValueError(f"Unsupported file format: {ext}")

    except ValueError:
        raise
    except Exception as exc:
        logger.exception("Unexpected extraction error for %s", file_path)
        raise ValueError("Unable to extract text from the uploaded file.") from exc

    import re
    # Clean up excessive newlines (e.g. \n\n\n -> \n\n)
    text = re.sub(r'\n{3,}', '\n\n', text)
    # Merge dangling prices like "Beef Biryani\n\nPKR 780" -> "Beef Biryani - PKR 780"
    text = re.sub(r'\n+(?=PKR\s*\d)', ' - ', text)

    return text


def split_text(text, chunk_size=1000, overlap=200):
    """Split text into overlapping chunks, respecting paragraphs and sentences."""
    if not text or not text.strip():
        return []

    text = text.replace('\r\n', '\n')
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = min(start + chunk_size, text_len)
        
        # Backtrack to a clean boundary if not at the end
        if end < text_len:
            last_newline = text.rfind('\n', start, end)
            if last_newline != -1 and last_newline > start + chunk_size // 2:
                end = last_newline + 1
            else:
                last_period = text.rfind('. ', start, end)
                if last_period != -1 and last_period > start + chunk_size // 2:
                    end = last_period + 2
                else:
                    last_space = text.rfind(' ', start, end)
                    if last_space != -1 and last_space > start:
                        end = last_space + 1
                        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
            
        if end >= text_len:
            break
            
        start = end - overlap
        if start < 0:
            start = 0
        else:
            next_newline = text.find('\n', start, end)
            if next_newline != -1 and next_newline < end - 50:
                start = next_newline + 1
            else:
                next_space = text.find(' ', start, end)
                if next_space != -1:
                    start = next_space + 1
                    
    return chunks